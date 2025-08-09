"""
PowerPoint Inconsistency Detection Tool

A comprehensive tool for detecting factual and logical inconsistencies
across PowerPoint presentations.

Installation:
    pip install python-pptx Pillow pytesseract google-generativeai python-dateutil

Usage:
    python ppt_analyzer.py presentation.pptx --api-key YOUR_GEMINI_API_KEY
"""

import argparse
import json
import logging
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, asdict
from io import BytesIO
import warnings
warnings.filterwarnings("ignore")

# -----------------------------------------------------------------------------
# Dependency Check (Not needed if you run the pip install commands first)
# -----------------------------------------------------------------------------
def check_and_install_dependencies():
    """Check for required dependencies and provide installation instructions."""
    missing_deps = []
    
    try:
        import pptx
    except ImportError:
        missing_deps.append("python-pptx")
    
    try:
        import PIL
    except ImportError:
        missing_deps.append("Pillow")
        
    try:
        import pytesseract
    except ImportError:
        missing_deps.append("pytesseract")
        
    try:
        import google.generativeai
    except ImportError:
        missing_deps.append("google-generativeai")
        
    try:
        import dateutil
    except ImportError:
        missing_deps.append("python-dateutil")
    
    if missing_deps:
        print(f"Missing required dependencies: {', '.join(missing_deps)}")
        print("\nTo install all dependencies, run:")
        print("pip install python-pptx Pillow pytesseract google-generativeai python-dateutil")
        print("\nFor OCR functionality, you also need to install Tesseract:")
        print("    Ubuntu/Debian: sudo apt-get install tesseract-ocr")
        print("    macOS: brew install tesseract")
        print("    Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki")
        return False
    
    return True

if not check_and_install_dependencies():
    sys.exit(1)

from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape
from PIL import Image
import pytesseract
import google.generativeai as genai
from dateutil import parser as date_parser
from google.colab import userdata # Colab specific import

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@dataclass
class SlideContent:
    """Container for extracted slide content."""
    slide_number: int
    text_content: str
    numerical_data: List[Dict[str, Any]]
    dates: List[Dict[str, Any]]
    images_text: str
    bullet_points: List[str]
    title: str
    
    def to_dict(self):
        """Convert to dictionary for JSON serialization."""
        return {
            'slide_number': self.slide_number,
            'title': self.title,
            'text_content': self.text_content[:200] + "..." if len(self.text_content) > 200 else self.text_content,
            'numerical_data_count': len(self.numerical_data),
            'dates_count': len(self.dates),
            'bullet_points_count': len(self.bullet_points),
            'has_images': bool(self.images_text.strip())
        }

@dataclass
class Inconsistency:
    """Container for detected inconsistencies."""
    type: str
    severity: str  # 'high', 'medium', 'low'
    description: str
    slides_involved: List[int]
    details: Dict[str, Any]
    confidence_score: float
    
    def to_dict(self):
        """Convert to dictionary for JSON serialization."""
        return {
            'type': self.type,
            'severity': self.severity,
            'description': self.description,
            'slides_involved': self.slides_involved,
            'details': self.details,
            'confidence_score': self.confidence_score
        }

class PowerPointAnalyzer:
    """Main analyzer class for PowerPoint inconsistency detection."""
    
    def __init__(self, api_key: str):
        """Initialize the analyzer with Gemini API key."""
        self.api_key = api_key
        self.slides_content: List[SlideContent] = []
        self.inconsistencies: List[Inconsistency] = []
        
        try:
            genai.configure(api_key=api_key)
            self.model = genai.GenerativeModel('gemini-2.5-flash')
            logger.info("Successfully configured Gemini AI")
        except Exception as e:
            logger.error(f"Failed to configure Gemini AI: {e}")
            raise
        
        self.number_patterns = [
            r'\$?\s*(\d{1,3}(?:,\d{3})*(?:\.\d+)?)\s*[MmBbKkTt]?(?:\s*(?:million|billion|thousand|M|B|K))?',
            r'(\d+(?:\.\d+)?)\s*%',
            r'(\d{1,3}(?:,\d{3})*(?:\.\d+)?)',
        ]
        
        self.date_patterns = [
            r'\b(\d{1,2}[-/]\d{1,2}[-/]\d{2,4})\b',
            r'\b(\d{4}[-/]\d{1,2}[-/]\d{1,2})\b',
            r'\b((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4})\b',
            r'\b(Q[1-4]\s+\d{4})\b',
            r'\b(FY\s*\d{4})\b',
        ]

    def extract_slide_content(self, presentation_path: str) -> List[SlideContent]:
        """Extract all content from PowerPoint slides."""
        try:
            if not os.path.exists(presentation_path):
                raise FileNotFoundError(f"Presentation file not found: {presentation_path}")
                
            prs = Presentation(presentation_path)
            logger.info(f"Processing presentation with {len(prs.slides)} slides")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                try:
                    content = self._extract_single_slide_content(slide, slide_idx)
                    self.slides_content.append(content)
                    logger.info(f"Processed slide {slide_idx}: {len(content.text_content)} characters, {len(content.numerical_data)} numbers")
                except Exception as e:
                    logger.warning(f"Error processing slide {slide_idx}: {e}")
                    self.slides_content.append(SlideContent(
                        slide_number=slide_idx,
                        text_content="",
                        numerical_data=[],
                        dates=[],
                        images_text="",
                        bullet_points=[],
                        title=f"Slide {slide_idx} (Processing Error)"
                    ))
            
            return self.slides_content
            
        except Exception as e:
            logger.error(f"Error processing presentation: {e}")
            raise

    def _extract_single_slide_content(self, slide, slide_number: int) -> SlideContent:
        """Extract content from a single slide."""
        text_content = ""
        bullet_points = []
        title = ""
        images_text = ""
        
        try:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        text_content += f"{text}\n"
                        
                        if not title and len(text) < 100 and not any(char in text for char in ['\n', '\t']):
                            title = text
                        
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.level > 0 or any(paragraph.text.strip().startswith(prefix) for prefix in ['•', '-', '*', '1.', '2.', '3.']):
                                    bullet_points.append(paragraph.text.strip())
                    
                    elif hasattr(shape, 'shapes'):
                        group_text = self._extract_from_group(shape)
                        if group_text.strip():
                            text_content += group_text
                    
                    elif shape.shape_type == 13:
                        try:
                            ocr_text = self._extract_text_from_image(shape)
                            if ocr_text.strip():
                                images_text += f"{ocr_text}\n"
                        except Exception as e:
                            logger.debug(f"OCR failed for image in slide {slide_number}: {e}")
                            
                except Exception as e:
                    logger.debug(f"Error processing shape in slide {slide_number}: {e}")
                    continue
            
        except Exception as e:
            logger.warning(f"Error extracting shapes from slide {slide_number}: {e}")
        
        all_text = text_content + images_text
        numerical_data = self._extract_numerical_data(all_text, slide_number)
        dates = self._extract_dates(all_text, slide_number)
        
        if not title:
            title = f"Slide {slide_number}"
        
        return SlideContent(
            slide_number=slide_number,
            text_content=text_content.strip(),
            numerical_data=numerical_data,
            dates=dates,
            images_text=images_text.strip(),
            bullet_points=bullet_points,
            title=title
        )

    def _extract_from_group(self, group_shape) -> str:
        """Extract text from grouped shapes."""
        text = ""
        try:
            for shape in group_shape.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text += f"{shape.text.strip()}\n"
                elif hasattr(shape, 'shapes'):
                    text += self._extract_from_group(shape)
        except Exception as e:
            logger.debug(f"Error extracting from group: {e}")
        return text

    def _extract_text_from_image(self, shape) -> str:
        """Extract text from image using OCR."""
        try:
            if hasattr(shape, 'image') and shape.image:
                image_stream = BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                
                custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz.,!?()[]{}:;-+=$%'
                text = pytesseract.image_to_string(image, config=custom_config)
                return text.strip()
        except Exception as e:
            logger.debug(f"OCR extraction failed: {e}")
        return ""

    def _extract_numerical_data(self, text: str, slide_number: int) -> List[Dict[str, Any]]:
        """Extract numerical data from text."""
        numerical_data = []
        
        for pattern in self.number_patterns:
            try:
                matches = re.finditer(pattern, text, re.IGNORECASE)
                for match in matches:
                    try:
                        full_match = match.group(0)
                        number_str = match.group(1)
                        
                        if re.match(r'^\d{4}$', number_str) and int(number_str) > 1900 and int(number_str) < 2100:
                            continue
                            
                        number = float(number_str.replace(',', ''))
                        
                        full_upper = full_match.upper()
                        if any(suffix in full_upper for suffix in ['K', 'THOUSAND']):
                            number *= 1000
                        elif any(suffix in full_upper for suffix in ['M', 'MILLION']):
                            number *= 1000000
                        elif any(suffix in full_upper for suffix in ['B', 'BILLION']):
                            number *= 1000000000
                        
                        data_type = 'number'
                        if '%' in full_match:
                            data_type = 'percentage'
                        elif '$' in full_match or any(word in text[max(0, match.start()-20):match.end()+20].lower() 
                                                     for word in ['revenue', 'cost', 'price', 'value', 'worth']):
                            data_type = 'currency'
                        
                        numerical_data.append({
                            'raw_text': full_match,
                            'value': number,
                            'context': text[max(0, match.start()-50):match.end()+50].strip(),
                            'slide': slide_number,
                            'type': data_type,
                            'position': match.start()
                        })
                        
                    except (ValueError, IndexError) as e:
                        logger.debug(f"Error parsing number {match.group(0)}: {e}")
                        continue
                        
            except Exception as e:
                logger.debug(f"Error with pattern {pattern}: {e}")
                continue
        
        return numerical_data

    def _extract_dates(self, text: str, slide_number: int) -> List[Dict[str, Any]]:
        """Extract dates from text."""
        dates = []
        
        for pattern in self.date_patterns:
            try:
                matches = re.finditer(pattern, text, re.IGNORECASE)
                for match in matches:
                    try:
                        date_str = match.group(1)
                        
                        parsed_date = None
                        date_type = 'standard'
                        
                        if date_str.upper().startswith('Q'):
                            date_type = 'quarter'
                        elif date_str.upper().startswith('FY'):
                            date_type = 'fiscal_year'
                        else:
                            try:
                                parsed_date = date_parser.parse(date_str, fuzzy=True)
                            except:
                                continue
                        
                        dates.append({
                            'raw_text': date_str,
                            'parsed_date': parsed_date.isoformat() if parsed_date else None,
                            'context': text[max(0, match.start()-30):match.end()+30].strip(),
                            'slide': slide_number,
                            'type': date_type
                        })
                        
                    except Exception as e:
                        logger.debug(f"Error parsing date {match.group(0)}: {e}")
                        continue
                        
            except Exception as e:
                logger.debug(f"Error with date pattern {pattern}: {e}")
                continue
        
        return dates

    def detect_inconsistencies(self) -> List[Inconsistency]:
        """Main method to detect all types of inconsistencies."""
        logger.info("Starting inconsistency detection...")
        
        try:
            self._detect_numerical_inconsistencies()
            self._detect_percentage_inconsistencies()
            self._detect_timeline_mismatches()
            self._detect_textual_contradictions()
            
            logger.info(f"Detected {len(self.inconsistencies)} potential inconsistencies")
            
        except Exception as e:
            logger.error(f"Error during inconsistency detection: {e}")
        
        return self.inconsistencies

    def _detect_numerical_inconsistencies(self):
        """Detect numerical inconsistencies across slides."""
        try:
            categories = {
                'revenue': [],
                'growth': [],
                'market': [],
                'costs': [],
                'users': []
            }
            
            for slide in self.slides_content:
                for num_data in slide.numerical_data:
                    context = num_data['context'].lower()
                    
                    if any(keyword in context for keyword in ['revenue', 'sales', 'income', 'earnings', 'turnover']):
                        categories['revenue'].append(num_data)
                    elif any(keyword in context for keyword in ['growth', 'increase', 'decrease', 'change', 'grew']):
                        categories['growth'].append(num_data)
                    elif any(keyword in context for keyword in ['market', 'share', 'size', 'cap', 'addressable']):
                        categories['market'].append(num_data)
                    elif any(keyword in context for keyword in ['cost', 'expense', 'spend', 'budget', 'investment']):
                        categories['costs'].append(num_data)
                    elif any(keyword in context for keyword in ['user', 'customer', 'client', 'subscriber']):
                        categories['users'].append(num_data)
            
            for category, data_list in categories.items():
                if len(data_list) >= 2:
                    self._check_numerical_group_consistency(category, data_list)
                    
        except Exception as e:
            logger.error(f"Error in numerical inconsistency detection: {e}")

    def _check_numerical_group_consistency(self, category: str, data_group: List[Dict]):
        """Check consistency within a group of numerical data."""
        try:
            for i, data1 in enumerate(data_group):
                for data2 in data_group[i+1:]:
                    if data1['slide'] != data2['slide']:
                        ratio = max(data1['value'], data2['value']) / max(min(data1['value'], data2['value']), 0.001)
                        
                        same_period = any(term in data1['context'].lower() and term in data2['context'].lower() 
                                         for term in ['q1', 'q2', 'q3', 'q4', '2023', '2024', 'current', 'this year'])
                        
                        if ratio > 2.0 and (same_period or category in ['revenue', 'market']):
                            severity = "high" if ratio > 10 else "medium" if ratio > 5 else "low"
                            confidence = min(0.9, 0.5 + (ratio - 2) * 0.1)
                            
                            self.inconsistencies.append(Inconsistency(
                                type="numerical_discrepancy",
                                severity=severity,
                                description=f"{category.title()} figures show {ratio:.1f}x discrepancy between slides",
                                slides_involved=[data1['slide'], data2['slide']],
                                details={
                                    'category': category,
                                    'value1': data1['value'],
                                    'value2': data2['value'],
                                    'context1': data1['context'][:100],
                                    'context2': data2['context'][:100],
                                    'ratio': round(ratio, 2)
                                },
                                confidence_score=confidence
                            ))
                            
        except Exception as e:
            logger.debug(f"Error checking group consistency for {category}: {e}")

    def _detect_percentage_inconsistencies(self):
        """Detect percentage calculation errors."""
        try:
            for slide in self.slides_content:
                percentages = [data for data in slide.numerical_data if data['type'] == 'percentage']
                
                if len(percentages) >= 2:
                    groups = []
                    for pct in percentages:
                        added_to_group = False
                        for group in groups:
                            if any(abs(pct['position'] - other['position']) < 200 for other in group):
                                group.append(pct)
                                added_to_group = True
                                break
                        if not added_to_group:
                            groups.append([pct])
                    
                    for group in groups:
                        if len(group) >= 2:
                            total = sum(p['value'] for p in group)
                            
                            if 80 <= total <= 120 and abs(total - 100) > 2:
                                self.inconsistencies.append(Inconsistency(
                                    type="percentage_calculation_error",
                                    severity="medium",
                                    description=f"Percentages total to {total:.1f}% instead of 100%",
                                    slides_involved=[slide.slide_number],
                                    details={
                                        'total_percentage': round(total, 1),
                                        'individual_percentages': [p['raw_text'] for p in group],
                                        'contexts': [p['context'][:50] for p in group]
                                    },
                                    confidence_score=0.8
                                ))
                                
        except Exception as e:
            logger.error(f"Error in percentage inconsistency detection: {e}")

    def _detect_timeline_mismatches(self):
        """Detect timeline and date inconsistencies."""
        try:
            dated_items = []
            for slide in self.slides_content:
                for date_data in slide.dates:
                    if date_data.get('parsed_date'):
                        try:
                            parsed = datetime.fromisoformat(date_data['parsed_date'])
                            dated_items.append({
                                'date': parsed,
                                'slide': slide.slide_number,
                                'context': date_data['context'].lower(),
                                'raw': date_data['raw_text']
                            })
                        except:
                            continue
            
            for i, item1 in enumerate(dated_items):
                for item2 in dated_items[i+1:]:
                    if item1['slide'] != item2['slide']:
                        context1, context2 = item1['context'], item2['context']
                        
                        is_forecast1 = any(word in context1 for word in ['forecast', 'projection', 'expected', 'planned', 'target'])
                        is_historical1 = any(word in context1 for word in ['actual', 'achieved', 'realized', 'historical'])
                        is_forecast2 = any(word in context2 for word in ['forecast', 'projection', 'expected', 'planned', 'target'])
                        is_historical2 = any(word in context2 for word in ['actual', 'achieved', 'realized', 'historical'])
                        
                        if (is_forecast1 and is_historical2 and item1['date'] < item2['date']) or \
                           (is_forecast2 and is_historical1 and item2['date'] < item1['date']):
                            self.inconsistencies.append(Inconsistency(
                                type="timeline_mismatch",
                                severity="medium",
                                description="Forecast date appears before historical date",
                                slides_involved=[item1['slide'], item2['slide']],
                                details={
                                    'date1': item1['raw'],
                                    'date2': item2['raw'],
                                    'context1': context1[:80],
                                    'context2': context2[:80]
                                },
                                confidence_score=0.7
                            ))
                            
        except Exception as e:
            logger.error(f"Error in timeline mismatch detection: {e}")

    def _detect_textual_contradictions(self):
        """Use Gemini AI to detect textual contradictions."""
        try:
            text_slides = [slide for slide in self.slides_content if len(slide.text_content.strip()) > 50]
            
            if len(text_slides) < 2:
                logger.info("Not enough text content for contradiction analysis")
                return
            
            slide_pairs = []
            for i, slide1 in enumerate(text_slides[:10]):
                for slide2 in text_slides[i+1:min(i+6, len(text_slides))]:
                    slide_pairs.append((slide1, slide2))
            
            batch_size = 3
            for i in range(0, min(len(slide_pairs), 9), batch_size):
                batch = slide_pairs[i:i+batch_size]
                try:
                    self._analyze_textual_batch(batch)
                except Exception as e:
                    logger.warning(f"Failed to analyze batch {i//batch_size + 1}: {e}")
                    continue
                    
        except Exception as e:
            logger.error(f"Error in textual contradiction detection: {e}")

    def _analyze_textual_batch(self, slide_pairs: List[Tuple[SlideContent, SlideContent]]):
        """Analyze a batch of slide pairs for textual contradictions."""
        try:
            prompt = self._create_contradiction_analysis_prompt(slide_pairs)
            
            response = self.model.generate_content(prompt)
            
            if response and hasattr(response, 'text'):
                self._parse_ai_contradiction_response(response.text, slide_pairs)
            else:
                logger.warning("Empty response from AI model")
                
        except Exception as e:
            logger.warning(f"AI analysis failed: {e}")

    def _create_contradiction_analysis_prompt(self, slide_pairs: List[Tuple[SlideContent, SlideContent]]) -> str:
        """Create prompt for AI contradiction analysis."""
        prompt = """Analyze the following PowerPoint slide pairs for logical contradictions, conflicting claims, or inconsistent statements.

Look for:
1. Contradictory factual claims (e.g., "market is growing" vs "market is declining")
2. Inconsistent strategic statements (e.g., "focus on cost reduction" vs "increase investment") 
3. Conflicting market assessments (e.g., "highly competitive" vs "few competitors")
4. Opposing conclusions or recommendations

Respond ONLY with valid JSON in this exact format:
{
  "contradictions": [
    {
      "slide1": <number>,
      "slide2": <number>,
      "type": "factual|strategic|market_assessment|conclusion",
      "severity": "high|medium|low",
      "description": "Brief description of contradiction",
      "evidence": {
        "statement1": "key contradicting phrase from slide 1",
        "statement2": "key contradicting phrase from slide 2"
      },
      "confidence": <0.6-0.9>
    }
  ]
}

Slide pairs:
"""
        
        for slide1, slide2 in slide_pairs[:3]:
            content1 = slide1.text_content[:300]
            content2 = slide2.text_content[:300]
            
            prompt += f"""
SLIDE {slide1.slide_number}: {slide1.title}
{content1}

SLIDE {slide2.slide_number}: {slide2.title}  
{content2}

---
"""
        
        return prompt

    def _parse_ai_contradiction_response(self, response_text: str, slide_pairs: List[Tuple[SlideContent, SlideContent]]):
        """Parse AI response for contradictions."""
        try:
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            
            if json_start >= 0 and json_end > json_start:
                json_str = response_text[json_start:json_end]
                data = json.loads(json_str)
                
                contradictions = data.get('contradictions', [])
                for contradiction in contradictions:
                    if isinstance(contradiction, dict):
                        slide1 = contradiction.get('slide1')
                        slide2 = contradiction.get('slide2')
                        
                        if slide1 and slide2:
                            self.inconsistencies.append(Inconsistency(
                                type="textual_contradiction",
                                severity=contradiction.get('severity', 'medium'),
                                description=contradiction.get('description', 'Textual contradiction detected'),
                                slides_involved=[slide1, slide2],
                                details=contradiction.get('evidence', {}),
                                confidence_score=contradiction.get('confidence', 0.6)
                            ))
                            
        except json.JSONDecodeError as e:
            logger.debug(f"Failed to parse AI response as JSON: {e}")
        except Exception as e:
            logger.debug(f"Error parsing AI response: {e}")

    def generate_report(self, output_file: Optional[str] = None) -> str:
        """Generate a comprehensive inconsistency report."""
        try:
            severity_order = {'high': 3, 'medium': 2, 'low': 1}
            sorted_inconsistencies = sorted(
                self.inconsistencies,
                key=lambda x: (severity_order.get(x.severity, 0), x.confidence_score),
                reverse=True
            )
            
            report_data = {
                'summary': {
                    'total_slides': len(self.slides_content),
                    'total_inconsistencies': len(self.inconsistencies),
                    'high_severity': len([i for i in self.inconsistencies if i.severity == 'high']),
                    'medium_severity': len([i for i in self.inconsistencies if i.severity == 'medium']),
                    'low_severity': len([i for i in self.inconsistencies if i.severity == 'low']),
                    'analysis_timestamp': datetime.now().isoformat()
                },
                'inconsistencies': [inc.to_dict() for inc in sorted_inconsistencies],
                'slide_summary': [slide.to_dict() for slide in self.slides_content]
            }
            
            if output_file:
                try:
                    with open(output_file, 'w', encoding='utf-8') as f:
                        json.dump(report_data, f, indent=2, ensure_ascii=False, default=str)
                    logger.info(f"Detailed report saved to {output_file}")
                except Exception as e:
                    logger.error(f"Failed to save report: {e}")
            
            return self._format_human_readable_report(report_data)
            
        except Exception as e:
            logger.error(f"Error generating report: {e}")
            return f"Error generating report: {e}"

    def _format_human_readable_report(self, report_data: Dict) -> str:
        """Format report for human reading."""
        try:
            summary = report_data['summary']
            inconsistencies = report_data['inconsistencies']
            
            output = f"""
PowerPoint Inconsistency Analysis
=================================

Analysis Summary
----------------
Total Slides Analyzed: {summary['total_slides']}
Inconsistencies Found: {summary['total_inconsistencies']}
  - High Severity: {summary['high_severity']}
  - Medium Severity: {summary['medium_severity']}
  - Low Severity: {summary['low_severity']}

"""
            
            if inconsistencies:
                output += "Detected Inconsistencies\n"
                output += "------------------------\n\n"
                
                for i, inc in enumerate(inconsistencies, 1):
                    slides_str = ", ".join(map(str, inc['slides_involved']))
                    
                    output += f"{i}. {inc['type'].replace('_', ' ').title()}\n"
                    output += f"    Severity: {inc['severity'].title()}\n"
                    output += f"    Slides: {slides_str}\n"
                    output += f"    Confidence: {inc['confidence_score']:.0%}\n"
                    output += f"    Description: {inc['description']}\n"
                    
                    details = inc.get('details', {})
                    if inc['type'] == 'numerical_discrepancy':
                        if 'value1' in details and 'value2' in details:
                            output += f"      Values: {details['value1']:,.0f} vs {details['value2']:,.0f}\n"
                            if 'ratio' in details:
                                output += f"      Ratio: {details['ratio']:.1f}x difference\n"
                    
                    elif inc['type'] == 'percentage_calculation_error':
                        if 'total_percentage' in details:
                            output += f"      Total: {details['total_percentage']}%\n"
                            if 'individual_percentages' in details:
                                output += f"      Parts: {', '.join(details['individual_percentages'])}\n"
                    
                    elif inc['type'] == 'textual_contradiction':
                        if 'statement1' in details and 'statement2' in details:
                            output += f"      Conflict:\n"
                            output += f"          \"{details['statement1'][:60]}...\"\n"
                            output += f"      vs  \"{details['statement2'][:60]}...\"\n"
                    
                    elif inc['type'] == 'timeline_mismatch':
                        if 'date1' in details and 'date2' in details:
                            output += f"      Dates: {details['date1']} vs {details['date2']}\n"
                    
                    output += "\n"
                
                output += "Recommendations\n"
                output += "---------------\n"
                
                high_count = summary['high_severity']
                medium_count = summary['medium_severity']
                
                if high_count > 0:
                    output += f"- URGENT: Review {high_count} high-severity inconsistencies immediately\n"
                    output += "  These likely represent significant data conflicts or errors\n\n"
                
                if medium_count > 0:
                    output += f"- IMPORTANT: Address {medium_count} medium-severity issues\n"
                    output += "  These may confuse stakeholders or indicate process gaps\n\n"
                
                types_found = set(inc['type'] for inc in inconsistencies)
                
                if 'numerical_discrepancy' in types_found:
                    output += "- Data Integrity: Verify all numerical data sources and calculations\n"
                if 'textual_contradiction' in types_found:
                    output += "- Content Review: Ensure consistent messaging across all slides\n"
                if 'percentage_calculation_error' in types_found:
                    output += "- Math Check: Review percentage calculations for accuracy\n"
                if 'timeline_mismatch' in types_found:
                    output += "- Timeline Audit: Verify all dates and forecasts for logical sequence\n"
                
            else:
                output += "No significant inconsistencies detected in your presentation.\n"
                output += "Your data appears to be well-aligned across all slides.\n"
            
            output += f"\nAnalysis completed at {summary['analysis_timestamp'][:19]}\n"
            
            return output
            
        except Exception as e:
            logger.error(f"Error formatting report: {e}")
            return f"Error formatting report: {e}"

def main():
    """Main function to run the analyzer."""
    # Use Colab's secrets manager to get the API key
    try:
        api_key = userdata.get('GEMINI_API_KEY')
        if not api_key:
            raise ValueError("GEMINI_API_KEY not found in Colab secrets. Please add it.")
    except Exception as e:
        print(f"Error retrieving API key from Colab secrets: {e}")
        sys.exit(1)

    # Define file path and output file directly in the code
    presentation_path = '/content/NoogatAssignment.pptx'
    output_file = 'report.json'
    verbose = False

    if not os.path.exists(presentation_path):
        print(f"Error: File not found: {presentation_path}")
        sys.exit(1)
    
    if not presentation_path.lower().endswith('.pptx'):
        print("Error: Please provide a .pptx file")
        sys.exit(1)
    
    print("Starting PowerPoint Inconsistency Analysis...")
    print(f"File: {presentation_path}")
    print("=" * 50)
    
    try:
        analyzer = PowerPointAnalyzer(api_key)
        
        print("Extracting slide content...")
        analyzer.extract_slide_content(presentation_path)
        
        print("Analyzing for inconsistencies...")
        inconsistencies = analyzer.detect_inconsistencies()
        
        print("Generating report...")
        
        readable_report = analyzer.generate_report(output_file)
        
        print("\n")
        print(readable_report)
        
        if len(inconsistencies) == 0:
            print("Great job! Your presentation appears to be consistent.")
        elif len(inconsistencies) <= 3:
            print(f"Analysis complete! Found {len(inconsistencies)} potential issues to review.")
        else:
            print(f"Analysis complete! Found {len(inconsistencies)} inconsistencies that need attention.")
        
        if output_file and os.path.exists(output_file):
            print(f"Detailed JSON report saved to: {output_file}")
        
    except KeyboardInterrupt:
        print("\nAnalysis cancelled by user")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Analysis failed: {e}")
        print(f"\nAnalysis failed: {e}")
        print("\nTroubleshooting tips:")
        print("1. Ensure your Gemini API key is valid and stored in Colab secrets as 'GEMINI_API_KEY'")
        print("2. Check that the PowerPoint file is not corrupted")
        print("3. Verify all dependencies are installed")
        print("4. Try with --verbose flag for more details")
        sys.exit(1)

if __name__ == "__main__":
    main()